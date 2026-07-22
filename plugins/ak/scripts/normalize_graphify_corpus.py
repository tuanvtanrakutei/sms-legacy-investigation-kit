#!/usr/bin/env python3
"""Build a deterministic, binary-free Graphify corpus for one app workspace."""

from __future__ import annotations

import argparse
import csv
import hashlib
import json
import shutil
import subprocess
import tempfile
from datetime import datetime, timezone
from pathlib import Path
from typing import Callable


BINARY_ACCESS = {".mdb", ".accdb", ".adp", ".laccdb", ".ldb"}
TEXT_SUFFIXES = {
    ".bas", ".cls", ".frm", ".vb", ".sql", ".txt", ".md", ".mdx", ".qmd", ".rst",
    ".json", ".yaml", ".yml", ".xml", ".html", ".htm", ".css", ".js", ".ts", ".tsx",
    ".jsx", ".py", ".ps1", ".vbs", ".bat", ".cmd", ".c", ".cpp", ".h", ".hpp",
    ".cs", ".java", ".go", ".rs", ".rb", ".php", ".sh", ".toml", ".ini", ".cfg",
}
TABULAR_SUFFIXES = {".csv", ".tsv"}
DOCUMENT_SUFFIXES = {".pdf", ".xlsx", ".xls", ".docx", ".pptx"}
IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"}
UNSUPPORTED_LEGACY = {".doc", ".ppt"}
FORBIDDEN_PARTS = {".git", "runs", "outputs", "evidence", "decisions", "secrets", "credentials"}
FORBIDDEN_NAMES = {".env", ".dsn"}
NORMALIZER_VERSION = "2.6.2"
MAX_ROWS = 10_000


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for chunk in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(chunk)
    return digest.hexdigest()


def decode_text(path: Path) -> tuple[str, str]:
    raw = path.read_bytes()
    encodings = ["utf-8-sig", "utf-16", "cp932", "shift_jis"]
    for encoding in encodings:
        try:
            return raw.decode(encoding), encoding.upper()
        except UnicodeDecodeError:
            continue
    raise UnicodeError("not valid UTF-8, UTF-16, CP932, or Shift-JIS")


def load_manifest(path: Path) -> dict:
    try:
        import yaml  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError("PyYAML is required for Graphify corpus normalization") from exc
    return yaml.safe_load(path.read_text(encoding="utf-8")) or {}


def output_dir_from(manifest: dict, app_root: Path) -> Path:
    relative = str(manifest.get("graphify", {}).get("output_dir", "graphify-out"))
    output = (app_root / relative).resolve()
    try:
        output.relative_to(app_root)
    except ValueError as exc:
        raise RuntimeError("graphify.output_dir must remain inside the app workspace") from exc
    return output


def declared_paths(manifest: dict) -> tuple[list[str], list[str]]:
    sources = manifest.get("sources", {}) or {}
    sql = sources.get("sql_server", {}) or {}
    japanese = sources.get("japanese_documents", {}) or {}
    values: list[str] = []
    for key in ("vba_exports", "screenshots", "reports", "sample_files", "app_documents"):
        values.extend(str(value) for value in sources.get(key, []) or [] if isinstance(value, str))
    values.extend(str(value) for value in sql.get("exported_paths", []) or [] if isinstance(value, str))
    if isinstance(japanese, dict):
        values.extend(str(value) for value in japanese.values() if isinstance(value, str) and value.strip())
    access = [str(item.get("path")) for item in sources.get("access_databases", []) or [] if isinstance(item, dict) and item.get("path")]
    return values, access


def component_paths(app_root: Path) -> list[str]:
    index = app_root / "extracted" / "component-index.json"
    if not index.is_file():
        return []
    try:
        data = json.loads(index.read_text(encoding="utf-8-sig"))
    except (OSError, ValueError):
        return []
    return [
        str(value)
        for component in data.get("components", [])
        if isinstance(component, dict)
        for value in component.get("source_paths", [])
        if isinstance(value, str)
    ]


def collect_sources(app_root: Path, manifest: dict) -> tuple[list[Path], list[Path], list[dict[str, str]]]:
    declared, access_declared = declared_paths(manifest)
    declared.extend(component_paths(app_root))
    declared.extend(["manifest.yaml", "extracted/component-index.json", "extracted/module-plan"])
    files: set[Path] = set()
    excluded_access: set[Path] = set()
    gaps: list[dict[str, str]] = []

    for relative in sorted(set(declared)):
        candidate = (app_root / relative).resolve()
        try:
            candidate.relative_to(app_root)
        except ValueError:
            gaps.append({"source_path": relative, "status": "OUTSIDE_WORKSPACE", "detail": "Path is outside the app workspace"})
            continue
        if candidate.is_file():
            files.add(candidate)
        elif candidate.is_dir():
            files.update(path for path in candidate.rglob("*") if path.is_file())
        elif relative not in {"extracted/component-index.json", "extracted/module-plan"}:
            gaps.append({"source_path": relative, "status": "MISSING", "detail": "Declared source does not exist"})

    for relative in access_declared:
        candidate = (app_root / relative).resolve()
        if candidate.is_file():
            excluded_access.add(candidate)
    access_root = app_root / "sources" / "access"
    if access_root.is_dir():
        excluded_access.update(path for path in access_root.rglob("*") if path.is_file() and path.suffix.lower() in BINARY_ACCESS)

    output = output_dir_from(manifest, app_root)
    safe_files: list[Path] = []
    for path in sorted(files):
        relative = path.relative_to(app_root)
        if output in path.parents:
            continue
        if any(part.lower() in FORBIDDEN_PARTS for part in relative.parts) or path.name.lower() in FORBIDDEN_NAMES or path.suffix.lower() == ".dsn":
            gaps.append({"source_path": relative.as_posix(), "status": "EXCLUDED_POLICY", "detail": "Secrets, credentials, run state, evidence, decisions, and outputs are not Graphify inputs"})
            continue
        if path.suffix.lower() in BINARY_ACCESS:
            excluded_access.add(path)
        else:
            safe_files.append(path)
    return safe_files, sorted(excluded_access), gaps


def markdown_table(rows: list[list[object]]) -> str:
    if not rows:
        return "_Empty table._\n"
    width = max(len(row) for row in rows)
    normalized = [[str(value if value is not None else "").replace("|", "\\|").replace("\n", "<br>") for value in row] + [""] * (width - len(row)) for row in rows]
    header = [f"Column {index + 1}" for index in range(width)]
    rendered = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    rendered.extend("| " + " | ".join(row) + " |" for row in normalized)
    return "\n".join(rendered) + "\n"


def normalize_delimited(path: Path) -> tuple[str, str, list[str]]:
    text, encoding = decode_text(path)
    delimiter = "\t" if path.suffix.lower() == ".tsv" else ","
    rows: list[list[object]] = []
    warnings: list[str] = []
    for index, row in enumerate(csv.reader(text.splitlines(), delimiter=delimiter)):
        if index >= MAX_ROWS:
            warnings.append(f"Truncated after {MAX_ROWS} rows")
            break
        rows.append(list(row))
    return markdown_table(rows), f"csv:{encoding}", warnings


def normalize_xlsx(path: Path) -> tuple[str, str, list[str]]:
    import openpyxl  # type: ignore[import-not-found]

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=False)
    sections: list[str] = []
    warnings: list[str] = []
    try:
        for sheet in workbook.worksheets:
            rows: list[list[object]] = []
            for index, row in enumerate(sheet.iter_rows(values_only=True)):
                if index >= MAX_ROWS:
                    warnings.append(f"Sheet {sheet.title!r} truncated after {MAX_ROWS} rows")
                    break
                rows.append(list(row))
            sections.append(f"## Sheet: {sheet.title}\n\n{markdown_table(rows)}")
    finally:
        workbook.close()
    return "\n".join(sections), f"openpyxl:{openpyxl.__version__}", warnings


def normalize_xls(path: Path) -> tuple[str, str, list[str]]:
    import xlrd  # type: ignore[import-not-found]

    workbook = xlrd.open_workbook(path, on_demand=True)
    sections: list[str] = []
    warnings: list[str] = []
    try:
        for sheet in workbook.sheets():
            count = min(sheet.nrows, MAX_ROWS)
            rows = [sheet.row_values(index) for index in range(count)]
            if sheet.nrows > count:
                warnings.append(f"Sheet {sheet.name!r} truncated after {MAX_ROWS} rows")
            sections.append(f"## Sheet: {sheet.name}\n\n{markdown_table(rows)}")
    finally:
        workbook.release_resources()
    return "\n".join(sections), f"xlrd:{xlrd.__version__}", warnings


def normalize_docx(path: Path) -> tuple[str, str, list[str]]:
    import docx  # type: ignore[import-not-found]

    document = docx.Document(path)
    sections = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    for index, table in enumerate(document.tables, 1):
        sections.append(f"## Table {index}\n\n" + markdown_table([[cell.text for cell in row.cells] for row in table.rows]))
    return "\n\n".join(sections), f"python-docx:{docx.__version__}", []


def normalize_pptx(path: Path) -> tuple[str, str, list[str]]:
    import pptx  # type: ignore[import-not-found]

    presentation = pptx.Presentation(path)
    sections: list[str] = []
    for index, slide in enumerate(presentation.slides, 1):
        texts = [shape.text for shape in slide.shapes if hasattr(shape, "text") and shape.text.strip()]
        sections.append(f"## Slide {index}\n\n" + "\n\n".join(texts))
    return "\n\n".join(sections), f"python-pptx:{pptx.__version__}", []


def tesseract_languages(executable: str) -> set[str]:
    result = subprocess.run([executable, "--list-langs"], check=False, capture_output=True, text=True)
    return {line.strip() for line in result.stdout.splitlines() if line.strip() and "available languages" not in line.lower()}


def ocr_images(images: list[Path]) -> tuple[str, str, list[str]]:
    executable = shutil.which("tesseract")
    if not executable:
        raise RuntimeError("OCR_REQUIRED: Tesseract executable is unavailable")
    languages = tesseract_languages(executable)
    selected = [value for value in ("jpn", "eng") if value in languages]
    if not selected:
        raise RuntimeError("OCR_REQUIRED: Tesseract has neither jpn nor eng language data")
    sections: list[str] = []
    for index, image in enumerate(images, 1):
        result = subprocess.run(
            [executable, str(image), "stdout", "-l", "+".join(selected)],
            check=False,
            capture_output=True,
            text=True,
            encoding="utf-8",
        )
        if result.returncode != 0:
            raise RuntimeError(f"OCR_FAILED: {result.stderr.strip()}")
        sections.append(f"## OCR page/image {index}\n\n{result.stdout.strip()}")
    return "\n\n".join(sections), f"tesseract:{'+'.join(selected)}", []


def normalize_pdf(path: Path) -> tuple[str, str, list[str]]:
    import pypdf  # type: ignore[import-not-found]

    reader = pypdf.PdfReader(path)
    page_texts = [(page.extract_text() or "").strip() for page in reader.pages]
    sections = [f"## Page {index}\n\n{text}" for index, text in enumerate(page_texts, 1)]
    combined = "\n\n".join(sections).strip()
    if any(page_texts):
        return combined, f"pypdf:{pypdf.__version__}", []

    try:
        import fitz  # type: ignore[import-not-found]
    except ImportError as exc:
        raise RuntimeError("OCR_REQUIRED: PDF has no text layer and PyMuPDF is unavailable") from exc
    with tempfile.TemporaryDirectory(prefix="ak-graphify-ocr-") as temp:
        images: list[Path] = []
        document = fitz.open(path)
        try:
            for index, page in enumerate(document):
                target = Path(temp) / f"page-{index + 1}.png"
                page.get_pixmap(matrix=fitz.Matrix(2, 2), alpha=False).save(target)
                images.append(target)
        finally:
            document.close()
        return ocr_images(images)


def normalize_image(path: Path) -> tuple[str, str, list[str]]:
    return ocr_images([path])


def normalize_source(path: Path) -> tuple[str, str, list[str], str]:
    suffix = path.suffix.lower()
    if suffix in TEXT_SUFFIXES:
        text, encoding = decode_text(path)
        destination_suffix = ".txt" if suffix in {".bas", ".cls", ".frm", ".vb"} else suffix
        return text, f"text:{encoding}", [], destination_suffix
    if suffix in TABULAR_SUFFIXES:
        text, parser, warnings = normalize_delimited(path)
        return text, parser, warnings, ".md"
    handlers: dict[str, Callable[[Path], tuple[str, str, list[str]]]] = {
        ".xlsx": normalize_xlsx,
        ".xls": normalize_xls,
        ".docx": normalize_docx,
        ".pptx": normalize_pptx,
        ".pdf": normalize_pdf,
    }
    if suffix in handlers:
        text, parser, warnings = handlers[suffix](path)
        return text, parser, warnings, ".md"
    if suffix in IMAGE_SUFFIXES:
        text, parser, warnings = normalize_image(path)
        return text, parser, warnings, ".md"
    if suffix in UNSUPPORTED_LEGACY:
        raise RuntimeError(f"CONVERSION_REQUIRED: legacy {suffix} requires conversion to DOCX/PPTX/PDF")
    raise RuntimeError(f"UNSUPPORTED_FORMAT: {suffix or '<no extension>'}")


def destination_for(corpus: Path, app_root: Path, source: Path, suffix: str) -> Path:
    relative = source.relative_to(app_root)
    target = corpus / "normalized" / relative
    return target.with_suffix(target.suffix + suffix if target.suffix.lower() != suffix else suffix)


def provenance_header(relative: str, source_hash: str, parser: str) -> str:
    return (
        "<!-- AK_GRAPHIFY_NORMALIZED\n"
        f"source_path: {relative}\nsource_sha256: {source_hash}\nparser: {parser}\n"
        f"normalizer_version: {NORMALIZER_VERSION}\n-->\n\n"
    )


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--app-root", required=True)
    parser.add_argument("--dry-run", action="store_true")
    parser.add_argument("--output")
    return parser.parse_args()


def main() -> int:
    args = parse_args()
    app_root = Path(args.app_root).expanduser().resolve()
    manifest_path = app_root / "manifest.yaml"
    if not manifest_path.is_file():
        raise SystemExit(f"Manifest not found: {manifest_path}")
    manifest = load_manifest(manifest_path)
    graph_root = output_dir_from(manifest, app_root)
    corpus = graph_root / "corpus"
    sources, excluded_access, initial_gaps = collect_sources(app_root, manifest)

    if args.dry_run:
        report = {
            "status": "PREFLIGHT_ONLY",
            "app_root": str(app_root),
            "corpus": str(corpus),
            "source_count": len(sources),
            "excluded_access_binaries": [path.relative_to(app_root).as_posix() for path in excluded_access],
            "planned_sources": [path.relative_to(app_root).as_posix() for path in sources],
        }
        print(json.dumps(report, ensure_ascii=False, indent=2))
        return 0

    if corpus.is_symlink():
        raise SystemExit(f"Refusing to replace symlinked Graphify corpus: {corpus}")
    if corpus.exists():
        shutil.rmtree(corpus)
    corpus.mkdir(parents=True)
    entries: list[dict[str, object]] = []
    total_words = 0
    for source in sources:
        relative = source.relative_to(app_root).as_posix()
        source_hash = sha256(source)
        try:
            text, parser_name, warnings, destination_suffix = normalize_source(source)
            target = destination_for(corpus, app_root, source, destination_suffix)
            target.parent.mkdir(parents=True, exist_ok=True)
            rendered = provenance_header(relative, source_hash, parser_name) + text.rstrip() + "\n"
            target.write_text(rendered, encoding="utf-8")
            output_hash = sha256(target)
            total_words += len(rendered.split())
            entries.append({
                "source_path": relative,
                "source_sha256": source_hash,
                "status": "NORMALIZED",
                "parser": parser_name,
                "output_path": target.relative_to(app_root).as_posix(),
                "output_sha256": output_hash,
                "warnings": warnings,
            })
        except (ImportError, OSError, RuntimeError, UnicodeError, ValueError) as exc:
            detail = str(exc)
            status = detail.split(":", 1)[0] if ":" in detail else "NORMALIZATION_FAILED"
            entries.append({"source_path": relative, "source_sha256": source_hash, "status": status, "detail": detail})

    for path in excluded_access:
        entries.append({
            "source_path": path.relative_to(app_root).as_posix(),
            "source_sha256": sha256(path),
            "status": "EXCLUDED_BINARY",
            "detail": "Access binaries and snapshots are never Graphify inputs",
        })
    entries.extend(initial_gaps)
    normalized = [entry for entry in entries if entry["status"] == "NORMALIZED"]
    fingerprint_input = [
        {key: entry.get(key) for key in ("source_path", "source_sha256", "output_path", "output_sha256", "parser", "status")}
        for entry in sorted(entries, key=lambda value: str(value.get("source_path")))
    ]
    fingerprint = hashlib.sha256(json.dumps(fingerprint_input, ensure_ascii=False, sort_keys=True).encode("utf-8")).hexdigest()
    gaps = [entry for entry in entries if entry["status"] not in {"NORMALIZED", "EXCLUDED_BINARY"}]
    report = {
        "schema_version": "2.1",
        "normalizer_version": NORMALIZER_VERSION,
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "status": "READY_WITH_GAPS" if normalized and gaps else ("READY" if normalized else "BLOCKED"),
        "app_root": str(app_root),
        "corpus_root": corpus.relative_to(app_root).as_posix(),
        "corpus_fingerprint": fingerprint,
        "corpus_file_count": len(normalized),
        "total_words": total_words,
        "excluded_access_binary_count": len(excluded_access),
        "binary_files_ingested": 0,
        "gap_count": len(gaps),
        "entries": entries,
    }
    graph_root.mkdir(parents=True, exist_ok=True)
    audit_path = graph_root / "CORPUS_AUDIT.json"
    audit_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    rendered = json.dumps(report, ensure_ascii=False, indent=2)
    print(rendered)
    if args.output:
        Path(args.output).expanduser().resolve().write_text(rendered + "\n", encoding="utf-8")
    return 0 if report["status"] != "BLOCKED" else 3


if __name__ == "__main__":
    raise SystemExit(main())
