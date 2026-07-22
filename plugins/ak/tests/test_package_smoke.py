from __future__ import annotations

import json
import subprocess
import sys
from pathlib import Path

import yaml


PACKAGE = Path(__file__).resolve().parents[1]
REPOSITORY = PACKAGE.parents[1]
SCRIPTS = PACKAGE / "scripts"
sys.path.insert(0, str(SCRIPTS))

from graphify_phase_gate import gate_status  # noqa: E402


def run_script(name: str, *args: str) -> subprocess.CompletedProcess[str]:
    result = subprocess.run(
        [sys.executable, str(SCRIPTS / name), *map(str, args)],
        cwd=PACKAGE,
        check=False,
        capture_output=True,
        text=True,
    )
    assert result.returncode == 0, f"{name} failed\nSTDOUT:\n{result.stdout}\nSTDERR:\n{result.stderr}"
    return result


def test_public_package_contract() -> None:
    import re

    package = json.loads((PACKAGE / "specifications" / "package.json").read_text(encoding="utf-8"))
    version = package["version"]
    assert re.fullmatch(r"\d+\.\d+\.\d+", version), version
    assert package["architecture_inspiration"]["dependency"] is False
    assert package["architecture_inspiration"]["vendored_code"] is False

    # package.json is the single source of truth; every manifest bump_version.py
    # touches must stay in lock-step so one release command is sufficient.
    for manifest in (".codex-plugin/plugin.json", ".claude-plugin/plugin.json"):
        assert json.loads((PACKAGE / manifest).read_text(encoding="utf-8"))["version"] == version
    marketplace = json.loads((REPOSITORY / ".claude-plugin" / "marketplace.json").read_text(encoding="utf-8"))
    assert marketplace["metadata"]["version"] == version
    assert all(entry["version"] == version for entry in marketplace["plugins"])

    run_script("validate_structure.py", "--package", str(PACKAGE), "--repository-root", str(REPOSITORY))


def test_public_yaml_has_unique_keys() -> None:
    class UniqueKeyLoader(yaml.SafeLoader):
        pass

    def construct_mapping(loader: UniqueKeyLoader, node: yaml.MappingNode, deep: bool = False) -> dict:
        mapping: dict = {}
        for key_node, value_node in node.value:
            key = loader.construct_object(key_node, deep=deep)
            assert key not in mapping, f"Duplicate YAML key {key!r} at line {key_node.start_mark.line + 1}"
            mapping[key] = loader.construct_object(value_node, deep=deep)
        return mapping

    UniqueKeyLoader.add_constructor(
        yaml.resolver.BaseResolver.DEFAULT_MAPPING_TAG,
        construct_mapping,
    )
    yaml_files = sorted((REPOSITORY / ".github").rglob("*.yml"))
    yaml_files += [REPOSITORY / "CITATION.cff", PACKAGE / "examples" / "minimal-app" / "manifest.yaml"]
    for path in yaml_files:
        yaml.load(path.read_text(encoding="utf-8"), Loader=UniqueKeyLoader)


def test_friendly_cli_entrypoint(tmp_path: Path) -> None:
    run_script("ak.py", "validate")
    run_script(
        "ak.py", "install",
        "--runtime", "generic",
        "--destination", str(tmp_path / "ak"),
        "--dry-run",
    )
    run_script(
        "ak.py", "install",
        "--runtime", "claude",
        "--project", str(tmp_path / "claude-project"),
        "--dry-run",
    )
    claude_project = tmp_path / "claude-project"
    run_script("ak.py", "install", "--runtime", "claude", "--project", str(claude_project))
    assert (claude_project / ".claude" / "ak-runtime").resolve() == PACKAGE.resolve()
    assert (claude_project / ".claude" / "skills" / "ak").resolve() == (PACKAGE / "skills" / "ak").resolve()
    run_script(
        "ak.py", "init",
        "--root", str(tmp_path),
        "--app-id", "T22",
        "--name-en", "Friendly CLI Test",
    )
    app = tmp_path / "T22"
    assert (app / "manifest.yaml").is_file()
    # Presentation output is optional and off by default.
    manifest_text = (app / "manifest.yaml").read_text(encoding="utf-8")
    assert "presentation_pptx: false" in manifest_text
    for token in ("required_before_phases: true", 'install_policy: "auto_managed"', 'refresh_policy: "before_each_phase"', 'corpus_policy: "binary_free_normalized"'):
        assert token in manifest_text
    run_script("ak.py", "preflight", "--app-root", str(app))
    planned_graph = run_script(
        "ak.py", "graphify", "prepare",
        "--app-root", str(app),
        "--phase", "1",
        "--runtime", "generic",
        "--dry-run",
    )
    assert json.loads(planned_graph.stdout)["status"] == "PREFLIGHT_ONLY"


def test_synthetic_module_aware_pipeline(tmp_path: Path) -> None:
    run_script(
        "init_app.py",
        "--root", str(tmp_path),
        "--app-id", "T21",
        "--name-en", "Synthetic Public Test",
        "--runtime", "generic",
    )
    app = tmp_path / "T21"

    access_file = app / "sources" / "access" / "synthetic.accdb"
    access_file.write_text("synthetic dry-run placeholder", encoding="utf-8")
    (app / "sources" / "vba" / "DemoForm.bas").write_text(
        'Attribute VB_Name = "DemoForm"\nSub Save_Click(): End Sub\n', encoding="utf-8"
    )
    ignored = app / "sources" / "sql" / "ignored.tmp"
    ignored.write_text("ignored", encoding="utf-8")

    compile_commands = app / "sources" / "documents" / "compile_commands.json"
    compile_commands.write_text(
        json.dumps([{
            "directory": "/synthetic/build",
            "file": "/synthetic/demo.cpp",
            "arguments": ["clang++", "--password", "secret-value", "-c", "/synthetic/demo.cpp"],
        }]),
        encoding="utf-8",
    )

    session = app / "extracted" / "access" / "DB1" / "20260715-000000"
    session.mkdir(parents=True)
    (session / "component-index.json").write_text(
        json.dumps({
            "schema_version": "2.1",
            "app_id": "T21",
            "generated_at": "2026-07-15T00:00:00+00:00",
            "components": [{
                "id": "DB1:form:DemoForm",
                "kind": "form",
                "name": "DemoForm",
                "container": "demo",
                "module_hint": "demo",
                "source_paths": ["forms/DemoForm.txt"],
                "depends_on": [],
                "metadata": {},
            }],
        }),
        encoding="utf-8",
    )

    dry_run = run_script(
        "extract_access.py",
        "--database", str(access_file),
        "--database-id", "DB1",
        "--session-id", "DRY-RUN",
        "--output-dir", str(app / "extracted" / "access"),
        "--dry-run",
    )
    assert json.loads(dry_run.stdout)["status"] == "PREFLIGHT_ONLY"

    normalized = app / "extracted" / "build-context" / "compile_commands.normalized.json"
    run_script("parse_compilation_database.py", "--input", str(compile_commands), "--output", str(normalized))
    normalized_text = normalized.read_text(encoding="utf-8")
    assert "secret-value" not in normalized_text
    assert "<REDACTED>" in normalized_text

    run_script("build_component_index.py", "--app-root", str(app))
    run_script(
        "build_module_plan.py",
        "--component-index", str(app / "extracted" / "component-index.json"),
        "--output-dir", str(app / "extracted" / "module-plan"),
    )
    run_script("create_run.py", "--app-root", str(app), "--runtime", "generic", "--run-id", "T21-PUBLIC-SMOKE")
    run = app / "runs" / "T21-PUBLIC-SMOKE"
    run_script("create_tasks.py", "--package", str(PACKAGE), "--run", str(run))

    inventory = json.loads((run / "source-inventory.json").read_text(encoding="utf-8"))
    assert any(item["relative_path"].endswith("ignored.tmp") for item in inventory["ignored"])
    tasks = [json.loads(path.read_text(encoding="utf-8")) for path in (run / "tasks").glob("*.json")]
    assert tasks
    assert any(task["module_targets"] for task in tasks)
    graph_tasks = [task for task in tasks if task["role"] == "graph_builder"]
    assert [task["phase_targets"] for task in graph_tasks] == [[phase] for phase in range(1, 7)]


def test_access_runtime_probe_reports_json() -> None:
    result = run_script("access_runtime.py")
    report = json.loads(result.stdout)
    for key in ("platform", "python_process_bitness", "registry", "powershell_candidates", "selected_host", "activation", "status"):
        assert key in report, f"access runtime report missing {key}"
    assert report["activation"]["status"] in {"NOT_REQUESTED", "READY", "NOT_INSTALLED", "REGISTERED_BUT_ACTIVATION_FAILED", "INSTALLED_BUT_BITNESS_MISMATCH", "NOT_FOUND"}
    assert report["selected_host"]["status"] in {"READY", "NOT_INSTALLED", "INSTALLED_BUT_BITNESS_MISMATCH", "NOT_FOUND"}


def test_extract_access_reports_runtime_block(tmp_path: Path) -> None:
    database = tmp_path / "runtime.accdb"
    database.write_text("synthetic dry-run placeholder", encoding="utf-8")

    discovered = run_script(
        "extract_access.py",
        "--database", str(database),
        "--database-id", "DBR",
        "--session-id", "DRY-RUNTIME",
        "--output-dir", str(tmp_path / "extracted"),
        "--dry-run",
    )
    plan = json.loads(discovered.stdout)
    assert plan["status"] == "PREFLIGHT_ONLY"
    runtime = plan["runtime"]
    assert runtime["runtime_check"] == "COMPLETED"
    assert runtime["runtime_tested"] is False
    assert "host" in runtime and "status" in runtime["host"]

    skipped = run_script(
        "extract_access.py",
        "--database", str(database),
        "--database-id", "DBR",
        "--session-id", "DRY-SKIP",
        "--output-dir", str(tmp_path / "extracted"),
        "--dry-run",
        "--skip-runtime-check",
    )
    skipped_runtime = json.loads(skipped.stdout)["runtime"]
    assert skipped_runtime["runtime_check"] == "SKIPPED"
    assert skipped_runtime["runtime_tested"] is False


def test_preflight_input_preconditions(tmp_path: Path) -> None:
    run_script(
        "init_app.py",
        "--root", str(tmp_path),
        "--app-id", "T24",
        "--name-en", "Preconditions Test",
        "--runtime", "generic",
    )
    app = tmp_path / "T24"
    manifest = app / "manifest.yaml"

    empty = json.loads(run_script("preflight.py", "--package", str(PACKAGE), "--manifest", str(manifest)).stdout)
    precond = empty["input_preconditions"]
    assert precond["mode"] == "none"
    assert set(precond["recommended_missing"]) == {"sources/vba", "sources/sql"}

    (app / "sources" / "vba" / "Form1.bas").write_text('Attribute VB_Name = "Form1"\n', encoding="utf-8")
    (app / "sources" / "sql" / "schema.sql").write_text("CREATE TABLE t(id int);\n", encoding="utf-8")
    exported = json.loads(run_script("preflight.py", "--package", str(PACKAGE), "--manifest", str(manifest)).stdout)
    assert exported["input_preconditions"]["mode"] == "export"
    assert exported["input_preconditions"]["recommended_missing"] == []

    access_db = app / "sources" / "access" / "T24.accdb"
    access_db.write_text("synthetic placeholder", encoding="utf-8")
    extract = json.loads(run_script("preflight.py", "--package", str(PACKAGE), "--manifest", str(manifest)).stdout)
    # VBA/SQL exports already present, so an unextracted Access binary makes it mixed.
    assert extract["input_preconditions"]["mode"] == "mixed"
    assert "runtime_status" in extract["input_preconditions"]


def test_nested_manifest_sources_drive_preflight_and_task_inputs(tmp_path: Path) -> None:
    run_script(
        "init_app.py",
        "--root", str(tmp_path),
        "--app-id", "T25",
        "--name-en", "Nested Source Test",
        "--runtime", "generic",
    )
    app = tmp_path / "T25"
    manifest = app / "manifest.yaml"
    text = manifest.read_text(encoding="utf-8")
    text = text.replace(
        'vba_exports: ["sources/vba"]',
        'vba_exports:\n    - "sources/T25_FRONTEND/vba"\n    - "sources/T25_DATA/vba"',
    )
    text = text.replace('exported_paths: ["sources/sql"]', "exported_paths: []")
    manifest.write_text(text, encoding="utf-8")
    frontend = app / "sources" / "T25_FRONTEND" / "vba"
    data = app / "sources" / "T25_DATA" / "vba"
    frontend.mkdir(parents=True)
    data.mkdir(parents=True)
    (frontend / "Form1.bas").write_text('Attribute VB_Name = "Form1"\n', encoding="utf-8")
    (data / "DataModule.bas").write_text('Attribute VB_Name = "DataModule"\n', encoding="utf-8")
    (app / "sources" / "access" / "T25.accdb").write_text("synthetic placeholder", encoding="utf-8")

    report = json.loads(run_script("preflight.py", "--package", str(PACKAGE), "--manifest", str(manifest)).stdout)
    preconditions = report["input_preconditions"]
    assert preconditions["mode"] == "mixed"
    assert preconditions["present"]["vba"] is True
    assert preconditions["present"]["sql"] is False
    assert preconditions["recommended_missing"] == []
    assert preconditions["present"]["present_paths"]["vba"] == [
        "sources/T25_FRONTEND/vba",
        "sources/T25_DATA/vba",
    ]

    run_script("create_run.py", "--app-root", str(app), "--runtime", "generic", "--run-id", "T25-NESTED")
    run = app / "runs" / "T25-NESTED"
    run_script("create_tasks.py", "--package", str(PACKAGE), "--run", str(run))
    tasks = [json.loads(path.read_text(encoding="utf-8")) for path in (run / "tasks").glob("*.json")]
    vba_task = next(task for task in tasks if task["role"] == "vba_ui")
    sql_task = next(task for task in tasks if task["role"] == "sql_data")
    assert "../../sources/T25_FRONTEND/vba" in vba_task["input_paths"]
    assert "../../sources/T25_DATA/vba" in vba_task["input_paths"]
    assert "../../sources/sql" not in sql_task["input_paths"]


def test_extract_ps1_declares_unique_safe_names() -> None:
    # Regression guard (runs everywhere): the safe-name function must derive a
    # deterministic hash from the original name so distinct non-ASCII objects do
    # not sanitize to the same filename and overwrite each other on disk.
    ps1 = (SCRIPTS / "extract_access.ps1").read_text(encoding="utf-8")
    assert "function Get-SafeName" in ps1
    assert "ComputeHash" in ps1


def test_extract_ps1_safe_names_are_unique_on_collide() -> None:
    import re
    import shutil

    powershell = shutil.which("powershell") or shutil.which("pwsh")
    if not powershell:
        import pytest

        pytest.skip("PowerShell not available on this platform")

    script = (SCRIPTS / "extract_access.ps1").as_posix()
    command = (
        "$ErrorActionPreference='Stop';"
        f"$c = Get-Content -Raw -LiteralPath '{script}';"
        "$m = [regex]::Match($c, '(?ms)^function Get-SafeName\\(.*?^\\}');"
        "if (-not $m.Success) { throw 'Get-SafeName not found' };"
        "Invoke-Expression $m.Value;"
        "$names = @('***','///','@@@','Form1','');"
        "($names | ForEach-Object { Get-SafeName $_ }) -join [char]10"
    )
    result = subprocess.run(
        [powershell, "-NoProfile", "-NonInteractive", "-Command", command],
        check=False,
        capture_output=True,
        text=True,
    )
    assert result.returncode == 0, result.stderr
    safe = [line.strip() for line in result.stdout.splitlines() if line.strip()]
    assert len(safe) == 5, safe
    assert len(set(safe)) == 5, f"safe names collide: {safe}"
    for name in safe:
        assert re.fullmatch(r"[A-Za-z0-9_.-]+", name), name


def test_vba_export_tool_present() -> None:
    bas = (PACKAGE / "tools" / "ExportAccessObjects.bas").read_text(encoding="utf-8")
    assert "Public Sub ExportAccessObjects" in bas
    assert "SaveAsText" in bas
    # Keeps original names, only de-duplicates on real collision (no lossy sanitizing).
    assert "UniquePath" in bas
    # Module name must differ from the Sub name, or calling it errors with
    # "Expected variable or procedure, not module".
    assert 'Attribute VB_Name = "modExportAccess"' in bas
    # SaveAsText output (system codepage / Shift-JIS) is transcoded to UTF-8 so
    # every export file is one consistent encoding.
    assert "shift_jis" in bas and 'Charset = "UTF-8"' in bas
    # Each object is exported independently; failures are recorded, not fatal.
    assert "AddSkip" in bas
    # System/temp/ImportErrors tables are excluded from the schema export.
    assert "IsSystemOrJunkTable" in bas


def test_recommended_optional_evidence_template() -> None:
    tpl = (PACKAGE / "templates" / "recommended-optional-evidence.md").read_text(encoding="utf-8")
    assert "{{APP_ID}}" in tpl
    assert "optional, non-blocking" in tpl
    for phase in ("Phase 4", "Phase 5", "Phase 2"):
        assert phase in tpl


def test_extract_access_blocks_before_snapshot_when_runtime_unavailable(tmp_path: Path) -> None:
    database = tmp_path / "blocked.accdb"
    database.write_text("synthetic placeholder", encoding="utf-8")
    out_dir = tmp_path / "extracted"
    # A non-existent PowerShell host forces NOT_FOUND without ever activating Access.
    result = subprocess.run(
        [
            sys.executable, str(SCRIPTS / "extract_access.py"),
            "--database", str(database),
            "--database-id", "DBX",
            "--session-id", "S1",
            "--output-dir", str(out_dir),
            "--execute",
            "--powershell", str(tmp_path / "missing" / "powershell.exe"),
        ],
        cwd=PACKAGE,
        check=False,
        capture_output=True,
        text=True,
    )
    assert result.returncode == 3, f"STDOUT:\n{result.stdout}\nSTDERR:\n{result.stderr}"
    plan = json.loads(result.stdout)
    assert plan["status"] == "BLOCKED"
    # The block must happen before any snapshot copy.
    assert not (out_dir / "DBX" / "S1" / "snapshot").exists()


def test_adopt_existing_workspace_preserves_files(tmp_path: Path) -> None:
    app = tmp_path / "T23"
    original = app / "docs" / "scope.md"
    original.parent.mkdir(parents=True)
    original.write_text("preserve", encoding="utf-8")

    refused = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS / "ak.py"),
            "init",
            "--app-root", str(app),
            "--app-id", "T23",
            "--name-en", "Adopted Existing App",
        ],
        cwd=PACKAGE,
        check=False,
        capture_output=True,
        text=True,
    )
    assert refused.returncode != 0
    assert not (app / "manifest.yaml").exists()

    run_script(
        "ak.py",
        "init",
        "--app-root", str(app),
        "--app-id", "T23",
        "--name-en", "Adopted Existing App",
        "--adopt-existing",
    )
    assert original.read_text(encoding="utf-8") == "preserve"
    assert (app / "manifest.yaml").is_file()


def test_graphify_runtime_dry_run_is_managed_and_non_mutating(tmp_path: Path) -> None:
    runtime_root = tmp_path / "managed-graphify"
    env = dict(__import__("os").environ)
    env["AK_GRAPHIFY_RUNTIME_ROOT"] = str(runtime_root)
    result = subprocess.run(
        [
            sys.executable,
            str(SCRIPTS / "graphify_runtime.py"),
            "ensure",
            "--runtime", "generic",
            "--dry-run",
        ],
        cwd=PACKAGE,
        env=env,
        check=False,
        capture_output=True,
        text=True,
    )
    assert result.returncode == 0, result.stderr
    report = json.loads(result.stdout)
    assert report["status"] == "INSTALL_PLANNED"
    assert report["isolated_from_system_python"] is True
    assert str(runtime_root) in report["runtime_root"]
    assert not runtime_root.exists()


def test_graphify_corpus_is_binary_free_and_cp932_safe(tmp_path: Path) -> None:
    run_script(
        "init_app.py",
        "--root", str(tmp_path),
        "--app-id", "T26",
        "--name-en", "Graphify Corpus Test",
    )
    app = tmp_path / "T26"
    vba = app / "sources" / "vba" / "受注処理.bas"
    vba.write_bytes('Attribute VB_Name = "受注処理"\nSub 実行(): End Sub\n'.encode("cp932"))
    (app / "sources" / "documents" / "rules.csv").write_text("項目,規則\n受注,必須\n", encoding="utf-8")
    (app / "sources" / "documents" / ".env").write_text("SECRET=do-not-ingest\n", encoding="utf-8")
    access = app / "sources" / "access" / "private.mdb"
    access.write_bytes(b"synthetic-access-binary")

    run_script("normalize_graphify_corpus.py", "--app-root", str(app))
    audit = json.loads((app / "graphify-out" / "CORPUS_AUDIT.json").read_text(encoding="utf-8"))
    assert audit["status"] in {"READY", "READY_WITH_GAPS"}
    assert audit["binary_files_ingested"] == 0
    assert audit["excluded_access_binary_count"] == 1
    assert all(not str(entry.get("output_path", "")).lower().endswith(".mdb") for entry in audit["entries"])
    assert next(entry for entry in audit["entries"] if entry.get("source_path", "").endswith("/.env"))["status"] == "EXCLUDED_POLICY"
    normalized_vba = next(
        app / entry["output_path"]
        for entry in audit["entries"]
        if entry.get("source_path", "").endswith("受注処理.bas")
    )
    assert "受注処理" in normalized_vba.read_text(encoding="utf-8")


def test_graphify_phase_gate_requires_fresh_graph_and_phase_receipt(tmp_path: Path) -> None:
    run_script(
        "init_app.py",
        "--root", str(tmp_path),
        "--app-id", "T27",
        "--name-en", "Graph Gate Test",
    )
    app = tmp_path / "T27"
    (app / "sources" / "vba" / "Form1.bas").write_text('Attribute VB_Name = "Form1"\n', encoding="utf-8")
    run_script("normalize_graphify_corpus.py", "--app-root", str(app))
    runtime = {"status": "READY", "requested_version": "0.9.18", "installed_version": "0.9.18", "graphify_executable": "synthetic"}
    missing = gate_status(app, 1, runtime)
    assert missing["reason"] == "GRAPH_BUILD_REQUIRED"

    root = app / "graphify-out"
    graph = root / "graph.json"
    graph.write_text(json.dumps({"nodes": [{"id": "Form1"}], "edges": []}), encoding="utf-8")
    audit = json.loads((root / "CORPUS_AUDIT.json").read_text(encoding="utf-8"))
    graph_hash = __import__("hashlib").sha256(graph.read_bytes()).hexdigest()
    state = {
        "corpus_fingerprint": audit["corpus_fingerprint"],
        "graph_sha256": graph_hash,
        "phase_queries": {
            "1": {
                "corpus_fingerprint": audit["corpus_fingerprint"],
                "output_path": "graphify-out/phase-context/phase-1.json",
            }
        },
    }
    (root / "GRAPH_STATE.json").write_text(json.dumps(state), encoding="utf-8")
    assert gate_status(app, 1, runtime)["status"] == "READY"
    assert gate_status(app, 2, runtime)["reason"] == "PHASE_QUERY_REQUIRED"
    (root / ".needs_update").write_text("semantic refresh pending", encoding="utf-8")
    assert gate_status(app, 1, runtime)["reason"] == "GRAPH_SEMANTIC_UPDATE_REQUIRED"
    (root / ".needs_update").unlink()

    (app / "sources" / "vba" / "Form1.bas").write_text('Attribute VB_Name = "Form1"\nSub Changed(): End Sub\n', encoding="utf-8")
    run_script("normalize_graphify_corpus.py", "--app-root", str(app))
    assert gate_status(app, 1, runtime)["reason"] == "GRAPH_UPDATE_REQUIRED"


def test_graphify_document_normalizers_cover_office_and_report_scanned_pdf(tmp_path: Path) -> None:
    import docx
    import openpyxl
    import pptx
    import pypdf

    run_script(
        "init_app.py",
        "--root", str(tmp_path),
        "--app-id", "T28",
        "--name-en", "Document Normalization Test",
    )
    app = tmp_path / "T28"
    documents = app / "sources" / "documents"

    workbook = openpyxl.Workbook()
    workbook.active.title = "業務規則"
    workbook.active.append(["項目", "規則"])
    workbook.active.append(["受注", "必須"])
    workbook.save(documents / "rules.xlsx")

    word = docx.Document()
    word.add_heading("運用手順", level=1)
    word.add_paragraph("受注確認を実行する。")
    word.save(documents / "manual.docx")

    slides = pptx.Presentation()
    slide = slides.slides.add_slide(slides.slide_layouts[1])
    slide.shapes.title.text = "業務フロー"
    slide.placeholders[1].text = "入力から出力まで"
    slides.save(documents / "flow.pptx")

    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=200, height=200)
    with (documents / "scan.pdf").open("wb") as handle:
        writer.write(handle)

    run_script("normalize_graphify_corpus.py", "--app-root", str(app))
    audit = json.loads((app / "graphify-out" / "CORPUS_AUDIT.json").read_text(encoding="utf-8"))
    statuses = {entry["source_path"]: entry["status"] for entry in audit["entries"]}
    assert statuses["sources/documents/rules.xlsx"] == "NORMALIZED"
    assert statuses["sources/documents/manual.docx"] == "NORMALIZED"
    assert statuses["sources/documents/flow.pptx"] == "NORMALIZED"
    assert statuses["sources/documents/scan.pdf"] in {"OCR_REQUIRED", "OCR_FAILED"}
